from pdf2image import convert_from_path
from subprocess import call, DEVNULL
from pptx import Presentation 
from gtts import gTTS
import os
import csv
import sys


def pptx2pdf(pptx, pdffolder='./'):
  """Convert pptx into pdf

  Args:
    pptx (str): pptx file
    pdffile (str): folder that the pdf is stored in

  Returns:
    str: pdf file path 
  """
  if not os.path.isdir(pdffolder):
    os.mkdir(pdffolder) 
  call(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', pdffolder, pptx], stdout=DEVNULL)
  pdfpath = os.path.join(pdffolder, os.path.basename(pptx).split('.')[0] + '.pdf')
  return pdfpath

def pdf2images(pdfpath, imgfolder='./', start=0, end=None):
  """Convert pdf into images

  Args:
    pdfpath (str): pdf file path
    imgfolder (str): image folder. Default is current folder
    start (int): start page of pdf
    end (int): end page pdf

  Returns:
    int: number of images stored
  """
  if not os.path.isdir(imgfolder):
    os.mkdir(imgfolder) 
  images = convert_from_path(pdfpath, thread_count=2, use_pdftocairo=True, size=(800, None), timeout=240)
  for index, image in enumerate(images[start:end]):
    image.save(os.path.join(imgfolder, str(index)+".jpg"))
  return len(images[start:end])


def gen_pdf_images(pptx):
  _tmp = "./_imgs"
  pdf = pptx2pdf(pptx, _tmp)
  size = pdf2images(pdf, _tmp)
  return (_tmp, size)

def gen_pptx_voices(pptx, lang):
  _tmp = "./_voices"
  if not os.path.isdir(_tmp):
    os.mkdir(_tmp) 

  slides = Presentation(pptx).slides
  size = len(slides)
  for i in range(size):
    note = slides[i].notes_slide
    text = note.notes_text_frame.text
    print(text)
    if text != "":
      mp3file = os.path.join(_tmp, "{}.mp3".format(i))
      if not os.path.isfile(mp3file):
        try:
          tts = gTTS(text, lang=lang)
          tts.save(mp3file)
        except Exception as e:
          print(e)
          print("{} has problem with voice!!!".format(text))
    else:
      call(['ffmpeg', '-f', 'lavfi', '-i', 'anullsrc=r=24000:cl=mono', '-t', '1', '-acodec', 'libmp3lame', '{}/{}.mp3'.format(_tmp, i)], stdout=DEVNULL)
  return (_tmp, size)


def merge_imgs_voices(imgs_folder, voices_folder, size):
  mp4_folder = './_mp4s'
  ts_folder = "./_tss"
  if not os.path.isdir(mp4_folder):
    os.mkdir(mp4_folder) 

  if not os.path.isdir(ts_folder):
    os.mkdir(ts_folder) 

  for i in range(size):
    jpg = "{}/{}.jpg".format(imgs_folder, i)
    mp3 = "{}/{}.mp3".format(voices_folder, i)
    mp4 = "{}/{}.mp4".format(mp4_folder, i)
    ts = "{}/{}.ts".format(ts_folder, i)
    if not os.path.isfile(ts):
      call(["ffmpeg", "-loop", "1", "-i", jpg, "-i", mp3, "-c:v", "libx264", "-tune", "stillimage", "-c:a", "aac", "-pix_fmt", "yuv420p", "-shortest", mp4], stdout=DEVNULL)
      call(["ffmpeg", "-y", "-i", mp4, "-c:v", "libx264", "-bsf:v", "h264_mp4toannexb", "-f", "mpegts", ts], stdout=DEVNULL) 

  return (ts_folder, size)

def merge_tss(ts_folder, size, video):
  tss = 'concat:' + '|'.join(["{}/{}.ts".format(ts_folder, i) for i in range(size)])
  call(["ffmpeg", "-y", "-f", "mpegts", "-i", tss, "-pix_fmt", "yuv420p", "-bsf:a", "aac_adtstoasc", video], stdout=DEVNULL)

def pptx2video(pptx, lang, video):
  imgs_folder, imgs_len = gen_pdf_images(pptx)
  voices_folder, voices_len = gen_pptx_voices(pptx, lang) 
  assert(imgs_len == voices_len)
  ts_folder, ts_len = merge_imgs_voices(imgs_folder, voices_folder, imgs_len)
  merge_tss(ts_folder, ts_len, video)
